import fitz
import pandas as pd
from docx import Document
import os
from google import genai
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import json

sentence_transformer_model = SentenceTransformer('all-MiniLM-L6-v2')

gc_client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
gc_model = "gemini-2.5-flash-lite"

def get_image_text(image_path):
    file = gc_client.files.upload(file=image_path)

    response = gc_client.models.generate_content(
        model=gc_model,
        contents=[
            {"role": "user", "parts": [
                {"file_data": {"file_uri": file.uri, "mime_type": file.mime_type}},
                {"text": "Extract all text from this document as it is even insert new line wherever needed."}
            ]}
        ]
    )
    return response.text

def get_audio_text(audio_path):
    file = gc_client.files.upload(file=audio_path)

    response = gc_client.models.generate_content(
        model=gc_model,
        contents=[
            {"role": "user", "parts": [
                {"file_data": {"file_uri": file.uri, "mime_type": file.mime_type}},
                {"text": "Extract exact text from given audio file and also distinguish between multiple speakers. In output there should be only text from audio and nothing else."}
            ]}
        ]
    )
    return response.text

def extract_text(file_path):
    text = ""
    if file_path.endswith(".pdf"):
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
    
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    
    elif file_path.endswith((".xlsx", ".xls")):
        df = pd.read_excel(file_path)
        df = df.fillna('')
        text = "".join(df.astype(str).values.flatten())
    
    elif file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
    
    elif file_path.endswith(".txt"):
        with open(file_path, "r") as f:
            text = f.read()
    
    elif file_path.endswith((".jpg",".jpeg",".png",".webp")):
        text = get_image_text(file_path)

    elif file_path.endswith((".ppt", ".pptx")):
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    
    elif file_path.endswith((".mp3",".wav")):
        text = get_audio_text(file_path)

    else:
        pass
    
    return text.strip()


def get_text():
    os.makedirs("./text_data",exist_ok=True)

    for file_name in os.listdir("./files"):
        if file_name == ".DS_Store":
            continue
        file_text = extract_text(f'./files/{file_name}')
        
        with open(f'./text_data/{file_name}.txt', 'w', encoding='utf-8') as f:
            f.write(file_text)
    
    return True

def get_top_filenames(query):
    query_embedding = sentence_transformer_model.encode([query])
    embeddings = np.load("./embedding_contents/embeddings.npy")
    sims = cosine_similarity(query_embedding, embeddings)[0]

    with open("./embedding_contents/embeded_files.json",'r') as f:
        file_list = json.load(f)
    file_sims = list(zip(file_list, sims))

    top_files = sorted(file_sims, key=lambda x: x[1], reverse=True)[:15]

    top_files_list = []
    for top_file in top_files:
        top_files_list.append(top_file[0])

    return top_files_list


def sync_embeddings():
    text_data_folder = "./text_data"
    embeddings_path = "./embedding_contents/embeddings.npy"
    embeddings_json = "./embedding_contents/embeded_files.json"

    text_data_files = []
    for filename in os.listdir(text_data_folder):
        text_data_files.append(filename[:-4])
    
    text_data = []
    for text_filename in os.listdir(text_data_folder):
        with open(f"./text_data/{text_filename}",'r', encoding='utf-8', errors='ignore') as f:
            text_data.append(f.read())
    
    embeddings = sentence_transformer_model.encode(text_data, show_progress_bar=True)
    np.save(embeddings_path, embeddings)

    os.makedirs(os.path.dirname(embeddings_json), exist_ok=True)

    with open(embeddings_json,'w') as f:
        json.dump(text_data_files,f,indent=4)
    
    return {'msg': 'All files embedded successfully!', 'files': text_data_files}